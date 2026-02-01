#!/usr/bin/env python3
"""
AI 助力公務效率提升 - PPTX 簡報生成器
執行: python3 generate_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 建立簡報
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 顏色定義
DARK_BG = RGBColor(15, 15, 26)
PURPLE = RGBColor(99, 102, 241)
LIGHT_PURPLE = RGBColor(165, 180, 252)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(148, 163, 184)
GREEN = RGBColor(52, 211, 153)

def add_title_slide(title, subtitle, meta=""):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 副標題
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = LIGHT_PURPLE
    p.alignment = PP_ALIGN.CENTER
    
    if meta:
        meta_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.333), Inches(0.6))
        tf = meta_box.text_frame
        p = tf.paragraphs[0]
        p.text = meta
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

def add_section_slide(num, title, time_info):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # 數字
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(2))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"0{num}"
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = RGBColor(60, 60, 100)
    p.alignment = PP_ALIGN.CENTER
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 時間
    time_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.333), Inches(0.6))
    tf = time_box.text_frame
    p = tf.paragraphs[0]
    p.text = time_info
    p.font.size = Pt(20)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(title, content_items):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 內容
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(24)
        p.font.color.rgb = GRAY
        p.space_after = Pt(12)

def add_step_slide(title, steps):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 步驟
    y_pos = 1.4
    for i, (step_title, step_desc) in enumerate(steps):
        # 數字圓圈
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y_pos), Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = PURPLE
        circle.line.fill.background()
        
        # 數字文字
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos + 0.1), Inches(0.6), Inches(0.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # 步驟標題
        step_title_box = slide.shapes.add_textbox(Inches(1.3), Inches(y_pos), Inches(11), Inches(0.5))
        tf = step_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = step_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # 步驟描述
        step_desc_box = slide.shapes.add_textbox(Inches(1.3), Inches(y_pos + 0.45), Inches(11), Inches(0.4))
        tf = step_desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = step_desc
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY
        
        y_pos += 1.2

def add_tool_slide(title, tools):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 工具卡片 (2x2 grid)
    positions = [(0.5, 1.4), (6.5, 1.4), (0.5, 4.2), (6.5, 4.2)]
    
    for i, (name, desc, url, tag) in enumerate(tools[:4]):
        x, y = positions[i]
        
        # 卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.8), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(30, 30, 50)
        card.line.color.rgb = PURPLE
        
        # 工具名稱
        name_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.3), Inches(5.2), Inches(0.5))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # 描述
        desc_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.8), Inches(5.2), Inches(0.8))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
        
        # URL
        url_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 1.5), Inches(5.2), Inches(0.4))
        tf = url_box.text_frame
        p = tf.paragraphs[0]
        p.text = url
        p.font.size = Pt(14)
        p.font.color.rgb = PURPLE
        
        # Tag
        tag_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 1.9), Inches(2), Inches(0.4))
        tf = tag_box.text_frame
        p = tf.paragraphs[0]
        p.text = tag
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = GREEN

def add_table_slide(title, headers, rows):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 表格
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.4), Inches(12.333), Inches(5)).table
    
    # 設定欄寬
    for col in table.columns:
        col.width = Inches(12.333 / num_cols)
    
    # 表頭
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = PURPLE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
    
    # 資料列
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(25, 25, 40)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(16)
            p.font.color.rgb = GRAY

def add_warning_slide():
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "⚠️ 公務使用注意事項"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    warnings = [
        ("🔒 資訊安全", "不要上傳機密文件或個人資料到 AI 工具"),
        ("✅ 人工審核", "AI 生成內容可能有錯誤，務必人工確認後才能使用"),
        ("©️ 著作權", "AI 生成圖片/文字的著作權問題仍有爭議，正式發布前請確認")
    ]
    
    y = 1.5
    for title, desc in warnings:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12.333), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(50, 20, 20)
        box.line.color.rgb = RGBColor(239, 68, 68)
        
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.2), Inches(11.7), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(252, 165, 165)
        
        desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.7), Inches(11.7), Inches(0.6))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY
        
        y += 1.8

# ===== 開始生成簡報 =====

# 封面
add_title_slide("🤖 AI 助力公務效率提升", "善用人工智慧工具，提升工作效能", "📍 文化部公務員培訓 ｜ ⏱️ 2 小時")

# 課程大綱
add_table_slide("📋 課程大綱", 
    ["段落", "主題", "時間"],
    [
        ["🔵 一", "Gemini：計劃構思與資料收集", "35 分鐘"],
        ["🟢 二", "AI 生成圖片與表格", "30 分鐘"],
        ["🟣 三", "AI 生成影片與字幕工具", "30 分鐘"],
        ["🔴 四", "AI 整理會議記錄", "20 分鐘"],
        ["⚪", "Q&A + 總結", "5 分鐘"]
    ])

# 第一部分
add_section_slide(1, "Google Gemini", "計劃構思與資料收集 ｜ 35 分鐘")

add_content_slide("什麼是 Google Gemini？", [
    "🧠 AI 智慧助手：Google 最新大型語言模型",
    "🌐 即時搜尋：整合 Google 搜尋能力",
    "📎 文件分析：可上傳 PDF、圖片進行分析",
    "💬 對話式介面：用自然語言溝通",
    "",
    "📍 網址：gemini.google.com",
    "💡 用 Google 帳號登入即可免費使用"
])

add_step_slide("實作步驟：開始使用 Gemini", [
    ("開啟 Gemini 網站", "在瀏覽器輸入 gemini.google.com"),
    ("登入 Google 帳號", "使用公務信箱或個人 Gmail 帳號登入"),
    ("開始對話", "在輸入框輸入你的問題或需求，按 Enter 送出"),
    ("持續追問", "根據回答繼續提問，讓 AI 幫你完善內容")
])

add_content_slide("實作：計劃構思", [
    "📌 情境：規劃一場文化藝術節活動",
    "",
    "💬 提示詞範例：",
    "「我要規劃一場為期三天的在地文化藝術節",
    "預算：500萬元，預期參與人數：5000人",
    "",
    "請幫我：",
    "1. 列出活動架構和時間表",
    "2. 建議的表演類型和攤位規劃", 
    "3. 需要注意的法規和申請事項",
    "4. 預算分配建議」",
    "",
    "💡 提示：具體描述需求，AI 才能給出精準回答"
])

add_content_slide("提示詞技巧", [
    "🎯 明確具體：說明目的、對象、預算、時間限制",
    "📋 條列需求：用 1、2、3 列出要 AI 做的事情",
    "🎭 設定角色：「請扮演資深活動策劃」",
    "📝 指定格式：「請用表格呈現」「限500字內」",
    "🔄 持續追問：「請更詳細說明第三點」",
    "📎 上傳文件：直接拖拉 PDF 或圖片讓 AI 分析"
])

# 第二部分
add_section_slide(2, "AI 生成圖片與表格", "視覺化工具介紹 ｜ 30 分鐘")

add_tool_slide("圖片生成工具推薦", [
    ("🎨 Canva AI", "中文介面、模板豐富", "canva.com", "⭐ 首推"),
    ("🖼️ Microsoft Designer", "整合 Office 365", "designer.microsoft.com", "公務適用"),
    ("🔥 Adobe Firefly", "商用安全、版權無慮", "firefly.adobe.com", "正式發布"),
    ("🆓 Bing Image Creator", "免費使用", "bing.com/create", "免費")
])

add_step_slide("實作：用 Canva 製作活動海報", [
    ("開啟 Canva → 選擇「海報」", "選擇適合的尺寸，如 A3 或社群貼文尺寸"),
    ("使用「AI 魔法設計」", "點擊左側「設計」→ 輸入活動主題關鍵字"),
    ("AI 生成圖片", "點擊「應用程式」→「AI 圖片產生器」→ 輸入描述"),
    ("下載成品", "右上角「分享」→「下載」→ 選擇 PNG 或 PDF")
])

add_content_slide("Gamma：AI 自動生成簡報", [
    "📍 網址：gamma.app",
    "",
    "🚀 使用步驟：",
    "1. 登入 gamma.app",
    "2. 點擊「Create new」",
    "3. 輸入簡報主題",
    "4. 選擇風格和頁數",
    "5. AI 自動生成！",
    "",
    "💬 範例：「文化資產保存政策簡報，包含現況分析、面臨挑戰、解決方案、預期成效」"
])

# 第三部分
add_section_slide(3, "AI 生成影片與字幕", "多媒體製作工具 ｜ 30 分鐘")

add_tool_slide("AI 影片生成工具", [
    ("🎬 Canva 影片", "模板豐富、操作直覺", "canva.com", "⭐ 入門推薦"),
    ("👤 Synthesia", "AI 虛擬主播", "synthesia.io", "專業級"),
    ("🎭 HeyGen", "AI 數位人、中文語音", "heygen.com", "專業級"),
    ("✈️ Runway", "文字/圖片轉影片", "runwayml.com", "進階")
])

add_tool_slide("字幕生成工具", [
    ("✂️ 剪映 CapCut", "免費、中文辨識超準", "capcut.com", "⭐ 強力推薦"),
    ("🌐 VEED.io", "線上工具、支援翻譯", "veed.io", "線上工具"),
    ("🎙️ 雅婷逐字稿", "工研院開發、台灣口音優化", "asr.yating.tw", "本土方案"),
    ("📺 YouTube 自動字幕", "上傳後自動產生", "youtube.com", "免費")
])

add_step_slide("實作：用剪映自動上字幕", [
    ("下載剪映（電腦版或手機版）", "官網下載，免費使用所有功能"),
    ("匯入影片檔案", "將錄好的影片拖入剪映時間軸"),
    ("點擊「文字」→「智能字幕」→「識別字幕」", "AI 自動辨識語音並生成字幕"),
    ("校對並匯出", "修正錯字，選擇匯出格式（影片或 SRT 字幕檔）")
])

# 第四部分
add_section_slide(4, "AI 整理會議記錄", "會議效率提升 ｜ 20 分鐘")

add_tool_slide("會議記錄工具", [
    ("🔥 Fireflies.ai", "自動加入會議錄音", "fireflies.ai", "⭐ 推薦"),
    ("📝 tl;dv", "錄製+轉錄+摘要", "tldv.io", "免費版夠用"),
    ("🏢 MS Copilot in Teams", "公務機關可能已有授權", "teams.microsoft.com", "公務適用"),
    ("🦦 Otter.ai", "即時轉錄+摘要", "otter.ai", "英文最強")
])

add_step_slide("本地方案（資安優先）", [
    ("會議錄音", "用手機或電腦錄音軟體"),
    ("語音轉文字", "用剪映或雅婷逐字稿轉成文字檔"),
    ("貼到 Gemini 整理", "請 AI 整理成：會議摘要、決議事項、待辦事項、下次議題"),
    ("人工審核", "確認內容正確後存檔")
])

# 總結
add_table_slide("📌 工具速查表",
    ["需求", "推薦工具"],
    [
        ["📝 計劃構思、資料整理", "Gemini、ChatGPT"],
        ["🎨 圖片生成", "Canva AI、Microsoft Designer"],
        ["📊 簡報製作", "Gamma、Canva"],
        ["🎬 影片製作", "Canva、剪映"],
        ["💬 字幕生成", "剪映、VEED.io"],
        ["📋 會議記錄", "剪映+Gemini、Fireflies"]
    ])

add_warning_slide()

add_title_slide("🙋 Q&A 時間", "有任何問題歡迎提問！", "感謝參與 ｜ 祝工作順利 🎉")

# 儲存
output_path = os.path.expanduser("~/Desktop/AI課程簡報.pptx")
prs.save(output_path)
print(f"✅ 簡報已儲存到：{output_path}")
print(f"📊 共 {len(prs.slides)} 頁")
